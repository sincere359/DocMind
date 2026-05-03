"""文档解析器：支持 PDF / TXT / Markdown / DOCX / XLSX / XLS / PPTX"""
from __future__ import annotations

import logging
from dataclasses import dataclass, field
from pathlib import Path
from typing import Iterator

from docmind.config import Config

logger = logging.getLogger(__name__)


@dataclass
class Chunk:
    """文本分块"""
    content: str
    metadata: dict = field(default_factory=dict)

    @property
    def source(self) -> str:
        return self.metadata.get("source", "unknown")

    @property
    def chunk_index(self) -> int:
        return self.metadata.get("chunk_index", 0)


class DocumentParser:
    """统一文档解析 + 分块"""

    SUPPORTED_EXTENSIONS = {".pdf", ".txt", ".md", ".docx", ".xlsx", ".xls", ".pptx"}

    def __init__(self, chunk_size: int | None = None, chunk_overlap: int | None = None):
        self.chunk_size = chunk_size or Config.CHUNK_SIZE
        self.chunk_overlap = chunk_overlap or Config.CHUNK_OVERLAP

    # ── 入口 ──

    def parse(self, file_path: str | Path) -> list[Chunk]:
        """解析文件，返回分块列表"""
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"文件不存在: {path}")

        suffix = path.suffix.lower()
        if suffix not in self.SUPPORTED_EXTENSIONS:
            raise ValueError(f"不支持的格式: {suffix}，支持: {self.SUPPORTED_EXTENSIONS}")

        raw_text = self._extract_text(path, suffix)
        if not raw_text.strip():
            raise ValueError(f"文件内容为空: {path.name}")

        chunks = self._split_text(raw_text, source=path.name)
        logger.info("解析完成: %s → %d 个分块", path.name, len(chunks))
        return chunks

    # ── 文本提取 ──

    def _extract_text(self, path: Path, suffix: str) -> str:
        extractors = {
            ".pdf": self._extract_pdf,
            ".txt": self._extract_txt,
            ".md": self._extract_txt,
            ".docx": self._extract_docx,
            ".xlsx": self._extract_excel,
            ".xls": self._extract_xls,
            ".pptx": self._extract_ppt,
        }
        return extractors[suffix](path)

    def _extract_pdf(self, path: Path) -> str:
        import pdfplumber

        pages = []
        with pdfplumber.open(path) as pdf:
            for i, page in enumerate(pdf.pages):
                text = page.extract_text() or ""
                if text.strip():
                    pages.append(f"[第{i+1}页]\n{text}")
        return "\n\n".join(pages)

    def _extract_txt(self, path: Path) -> str:
        return path.read_text(encoding="utf-8", errors="ignore")

    def _extract_docx(self, path: Path) -> str:
        from docx import Document

        doc = Document(str(path))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n\n".join(paragraphs)

    def _extract_excel(self, path: Path) -> str:
        """提取 Excel (.xlsx) 文件文本，逐 sheet 逐行提取"""
        from openpyxl import load_workbook

        wb = load_workbook(str(path), read_only=True, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c).strip() for c in row if c is not None]
                if cells:
                    rows.append(" | ".join(cells))
            if rows:
                parts.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(parts)

    def _extract_xls(self, path: Path) -> str:
        """提取旧版 Excel (.xls) 文件文本，逐 sheet 逐行提取"""
        import xlrd

        wb = xlrd.open_workbook(str(path))
        parts = []
        for sheet in wb.sheets():
            rows = []
            for row_idx in range(sheet.nrows):
                cells = [str(sheet.cell_value(row_idx, col_idx)).strip()
                         for col_idx in range(sheet.ncols)
                         if str(sheet.cell_value(row_idx, col_idx)).strip()]
                if cells:
                    rows.append(" | ".join(cells))
            if rows:
                parts.append(f"[Sheet: {sheet.name}]\n" + "\n".join(rows))
        return "\n\n".join(parts)

    def _extract_ppt(self, path: Path) -> str:
        """提取 PowerPoint (.pptx) 文件文本，逐 slide 提取"""
        from pptx import Presentation

        prs = Presentation(str(path))
        parts = []
        for i, slide in enumerate(prs.slides):
            texts = []
            if slide.shapes.title and slide.shapes.title.text.strip():
                texts.append(f"标题: {slide.shapes.title.text.strip()}")
            for shape in slide.shapes:
                if shape.has_text_frame and shape != slide.shapes.title:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            texts.append(text)
            if texts:
                parts.append(f"[幻灯片 {i+1}]\n" + "\n".join(texts))
        return "\n\n".join(parts)

    # ── 分块策略 ──

    def _split_text(self, text: str, source: str = "") -> list[Chunk]:
        """滑动窗口分块，保留上下文重叠"""
        chunks = []
        text_len = len(text)
        start = 0
        idx = 0

        while start < text_len:
            end = start + self.chunk_size
            chunk_text = text[start:end].strip()

            if chunk_text:
                # 尝试在句子边界处切分
                if end < text_len:
                    last_period = max(
                        chunk_text.rfind("。"),
                        chunk_text.rfind("！"),
                        chunk_text.rfind("？"),
                        chunk_text.rfind("\n"),
                    )
                    if last_period > self.chunk_size // 2:
                        chunk_text = chunk_text[: last_period + 1].strip()
                        end = start + last_period + 1

                chunks.append(Chunk(
                    content=chunk_text,
                    metadata={
                        "source": source,
                        "chunk_index": idx,
                        "start_char": start,
                        "end_char": end,
                    },
                ))
                idx += 1

            start = end - self.chunk_overlap
            if start <= end - self.chunk_size:
                start = end  # 防止死循环

        return chunks


def parse_document(file_path: str | Path) -> list[Chunk]:
    """便捷函数：解析单个文档"""
    return DocumentParser().parse(file_path)


def parse_documents(file_paths: list[str | Path]) -> list[Chunk]:
    """便捷函数：批量解析多个文档"""
    parser = DocumentParser()
    all_chunks = []
    for fp in file_paths:
        try:
            chunks = parser.parse(fp)
            all_chunks.extend(chunks)
        except Exception as e:
            logger.error("解析失败 %s: %s", fp, e)
    return all_chunks
